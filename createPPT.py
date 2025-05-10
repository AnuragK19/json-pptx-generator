from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import json
import requests
import os
from PIL import Image
import io

# Load JSON file
with open("sample_slides.json", "r") as f:
    slides_data = json.load(f)

# Load the template
prs = Presentation("custom_2x2_template.potx")

# Helper function to convert hex color to RGB
def hex_to_rgb(hex_color):
    hex_color = hex_color.replace("0x", "")
    return RGBColor(int(hex_color[2:4], 16), int(hex_color[4:6], 16), int(hex_color[6:8], 16))

# Helper function to map alignment
def get_alignment(alignment_x):
    if alignment_x == -1:
        return PP_ALIGN.LEFT
    elif alignment_x == 1:
        return PP_ALIGN.RIGHT
    else:
        return PP_ALIGN.CENTER

# Helper function to download and validate image
def download_image(url, filename):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        content = response.content
        image = Image.open(io.BytesIO(content))
        image.verify()
        image.close()
        with open(filename, "wb") as f:
            f.write(content)
        return filename
    except (requests.RequestException, Image.UnidentifiedImageError, Exception) as e:
        print(f"Failed to download or validate image from {url}: {str(e)}")
        return None

# Process each slide
for slide_idx, slide_data in enumerate(slides_data):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set a black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # Use the corresponding template slide
    template_slide = prs.slides[slide_idx]
    print(f"Processing slide {slide_idx + 1} with template shapes: {[shape.name for shape in template_slide.shapes]}")

    # Copy shapes from the template slide
    for template_shape in template_slide.shapes:
        if template_shape.has_text_frame:
            new_shape = slide.shapes.add_textbox(
                left=template_shape.left,
                top=template_shape.top,
                width=template_shape.width,
                height=template_shape.height
            )
            new_shape.name = template_shape.name
            new_shape.text = template_shape.text
            new_shape.text_frame.word_wrap = True
        else:
            new_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left=template_shape.left,
                top=template_shape.top,
                width=template_shape.width,
                height=template_shape.height
            )
            new_shape.name = template_shape.name

    print(f"Shapes in slide {slide_idx + 1} after copying: {[shape.name for shape in slide.shapes]}")

    # Process sections
    for idx, section in enumerate(slide_data["sections"]):
        placeholder_name = section.get("placeholderName")
        shape = None
        for s in slide.shapes:
            if s.name == placeholder_name:
                shape = s
                break

        if not shape:
            print(f"Shape {placeholder_name} not found in slide {slide_idx + 1}!")
            continue

        print(f"Populating {placeholder_name} in slide {slide_idx + 1}")

        if section["type"] == "content_text":
            text_frame = shape.text_frame
            text_frame.clear()
            
            # Split the text into paragraphs
            paragraphs = section["data"].replace("\r\n", "\n").split("\n")
            
            for para_text in paragraphs:
                # Add a new paragraph
                p = text_frame.add_paragraph()
                
                # Check if the paragraph should be a bullet point
                if para_text.strip().startswith("- "):
                    p.text = para_text.strip()[2:]  # Remove the "- " prefix
                    p.level = 0  # Bullet level (0 for top-level bullets)
                    p.bullet = True  # Enable bullet
                else:
                    p.text = para_text.strip()
                
                # Apply paragraph-level formatting
                p.alignment = get_alignment(section["alignmentX"])
                for run in p.runs:
                    run.font.size = Pt(section["fontSize"])
                    run.font.bold = section["isBold"]
                    run.font.color.rgb = hex_to_rgb(section["colorHex"])
            
            text_frame.word_wrap = True

        elif section["type"] == "image":
            image_url = section.get("imageUrl")
            if image_url:
                image_filename = f"temp_image_{slide_idx}_{idx}.png"
                image_path = download_image(image_url, image_filename)
                if image_path:
                    try:
                        slide.shapes.add_picture(
                            image_path,
                            left=shape.left,
                            top=shape.top,
                            width=shape.width,
                            height=shape.height
                        )
                        shape._element.getparent().remove(shape._element)
                    finally:
                        os.remove(image_path)
                else:
                    print(f"Skipping image for {placeholder_name} due to download failure in slide {slide_idx + 1}")
                    shape.text = f"Image failed to load: {image_url}"

    # Add speaker notes
    if slide_data["speakerNotes"]:
        slide.notes_slide.notes_text_frame.text = slide_data["speakerNotes"]

# Save the presentation
prs.save("output_presentation.pptx")
print("PowerPoint presentation created successfully!")