from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import requests
import os
from PIL import Image
import io
import tempfile

app = Flask(__name__)

# Initialize presentation for template
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # Blank layout

# Helper function to add a text box with formatting
def add_text_box(slide, name, left, top, width, height, text, font_size, is_bold=False, alignment=PP_ALIGN.LEFT):
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_box.name = name
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = is_bold

# Helper function to add an image placeholder
def add_image_placeholder(slide, name, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.name = name
    shape.text = f"{name} Placeholder"
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(14)

# Layout 1: Title Slide
def add_layout_1_title_slide():
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, "TitleSlideTitle", Inches(1), Inches(1), Inches(11.33), Inches(2), 
                 "TitleSlide Title Placeholder", 44, True, PP_ALIGN.CENTER)
    add_text_box(slide, "TitleSlideSubtitle", Inches(1), Inches(3.5), Inches(11.33), Inches(1), 
                 "TitleSlide Subtitle Placeholder", 24, False, PP_ALIGN.CENTER)
    add_image_placeholder(slide, "TitleSlideLogo", Inches(5.665), Inches(5.5), Inches(2), Inches(1))

# Layout 2: 2x2 Grid Layout
def add_layout_2_2x2_grid():
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, "GridTopLeftTitle", Inches(0.5), Inches(0.3), Inches(6), Inches(0.5), 
                 "GridTopLeftTitle Placeholder", 18, True)
    add_text_box(slide, "GridTopLeftContent", Inches(0.5), Inches(0.9), Inches(6), Inches(2.6), 
                 "GridTopLeftContent Placeholder", 14)
    add_image_placeholder(slide, "GridTopRightImage", Inches(6.83), Inches(0.5), Inches(6), Inches(3))
    add_image_placeholder(slide, "GridBottomLeftImage", Inches(0.5), Inches(4), Inches(6), Inches(3))
    add_text_box(slide, "GridBottomRightTitle", Inches(6.83), Inches(3.8), Inches(6), Inches(0.5), 
                 "GridBottomRightTitle Placeholder", 18, True)
    add_text_box(slide, "GridBottomRightContent", Inches(6.83), Inches(4.4), Inches(6), Inches(2.6), 
                 "GridBottomRightContent Placeholder", 14)

# Layout 3: Image with Caption and Content
def add_layout_3_image_caption_content():
    slide = prs.slides.add_slide(blank_layout)
    add_image_placeholder(slide, "ImageCaptionLeftImage", Inches(0.5), Inches(0.5), Inches(6), Inches(5.5))
    add_text_box(slide, "ImageCaptionLeftCaption", Inches(0.5), Inches(6.1), Inches(6), Inches(1), 
                 "ImageCaptionLeftCaption Placeholder", 12)
    add_text_box(slide, "ImageCaptionRightTitle", Inches(6.83), Inches(0.5), Inches(6), Inches(0.5), 
                 "ImageCaptionRightTitle Placeholder", 18, True)
    add_text_box(slide, "ImageCaptionRightContent", Inches(6.83), Inches(1.1), Inches(6), Inches(6.1), 
                 "ImageCaptionRightContent Placeholder", 14)

# Layout 4: Image Left, Two Sections Right
def add_layout_4_image_two_sections():
    slide = prs.slides.add_slide(blank_layout)
    add_image_placeholder(slide, "ImageTwoSectionsLeftImage", Inches(0.5), Inches(0.5), Inches(6), Inches(6.5))
    add_text_box(slide, "ImageTwoSectionsRightTitle1", Inches(6.83), Inches(0.5), Inches(6), Inches(0.5), 
                 "ImageTwoSectionsRightTitle1 Placeholder", 18, True)
    add_text_box(slide, "ImageTwoSectionsRightContent1", Inches(6.83), Inches(1.1), Inches(6), Inches(2.9), 
                 "ImageTwoSectionsRightContent1 Placeholder", 14)
    add_text_box(slide, "ImageTwoSectionsRightTitle2", Inches(6.83), Inches(4.2), Inches(6), Inches(0.5), 
                 "ImageTwoSectionsRightTitle2 Placeholder", 18, True)
    add_text_box(slide, "ImageTwoSectionsRightContent2", Inches(6.83), Inches(4.8), Inches(6), Inches(2.2), 
                 "ImageTwoSectionsRightContent2 Placeholder", 14)

# Layout 5: Full Slide Image with Overlay Text
def add_layout_5_full_image_overlay():
    slide = prs.slides.add_slide(blank_layout)
    add_image_placeholder(slide, "FullImageOverlayImage", Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    add_text_box(slide, "FullImageOverlayText", Inches(1), Inches(3), Inches(11.33), Inches(1.5), 
                 "FullImageOverlayText Placeholder", 32, True, PP_ALIGN.CENTER)

# Layout 6: Three-Column Layout
def add_layout_6_three_columns():
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, "ThreeColumnsTitle1", Inches(0.5), Inches(0.5), Inches(4), Inches(0.5), 
                 "ThreeColumnsTitle1 Placeholder", 18, True)
    add_text_box(slide, "ThreeColumnsContent1", Inches(0.5), Inches(1.1), Inches(4), Inches(6.1), 
                 "ThreeColumnsContent1 Placeholder", 14)
    add_text_box(slide, "ThreeColumnsTitle2", Inches(4.66), Inches(0.5), Inches(4), Inches(0.5), 
                 "ThreeColumnsTitle2 Placeholder", 18, True)
    add_text_box(slide, "ThreeColumnsContent2", Inches(4.66), Inches(1.1), Inches(4), Inches(6.1), 
                 "ThreeColumnsContent2 Placeholder", 14)
    add_text_box(slide, "ThreeColumnsTitle3", Inches(8.83), Inches(0.5), Inches(4), Inches(0.5), 
                 "ThreeColumnsTitle3 Placeholder", 18, True)
    add_text_box(slide, "ThreeColumnsContent3", Inches(8.83), Inches(1.1), Inches(4), Inches(6.1), 
                 "ThreeColumnsContent3 Placeholder", 14)

# Layout 7: Title with Two Images Below
def add_layout_7_title_two_images():
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, "TitleTwoImagesTitle", Inches(1), Inches(0.5), Inches(11.33), Inches(1), 
                 "TitleTwoImagesTitle Placeholder", 32, True, PP_ALIGN.CENTER)
    add_image_placeholder(slide, "TitleTwoImagesLeftImage", Inches(0.5), Inches(2), Inches(6), Inches(4))
    add_image_placeholder(slide, "TitleTwoImagesRightImage", Inches(6.83), Inches(2), Inches(6), Inches(4))

# Layout 8: Comparison Layout
def add_layout_8_comparison():
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, "ComparisonLeftTitle", Inches(0.5), Inches(0.5), Inches(6), Inches(0.5), 
                 "ComparisonLeftTitle Placeholder", 18, True)
    add_text_box(slide, "ComparisonLeftContent", Inches(0.5), Inches(1.1), Inches(6), Inches(6.1), 
                 "ComparisonLeftContent Placeholder", 14)
    add_text_box(slide, "ComparisonRightTitle", Inches(6.83), Inches(0.5), Inches(6), Inches(0.5), 
                 "ComparisonRightTitle Placeholder", 18, True)
    add_text_box(slide, "ComparisonRightContent", Inches(6.83), Inches(1.1), Inches(6), Inches(6.1), 
                 "ComparisonRightContent Placeholder", 14)

# Layout 9: Quote Slide
def add_layout_9_quote():
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, "QuoteText", Inches(1), Inches(1.5), Inches(11.33), Inches(3), 
                 "QuoteText Placeholder", 28, False, PP_ALIGN.CENTER)
    add_text_box(slide, "QuoteAttribution", Inches(1), Inches(5), Inches(11.33), Inches(1), 
                 "QuoteAttribution Placeholder", 18, False, PP_ALIGN.CENTER)

# Layout 10: Timeline Layout
def add_layout_10_timeline():
    slide = prs.slides.add_slide(blank_layout)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.5), Inches(11.33), Inches(0.1))
    shape.name = "TimelineLine"
    for i in range(4):
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1 + i*3.6), Inches(3.4), Inches(0.3), Inches(0.3))
        marker.name = f"TimelineMarker{i+1}"
        add_text_box(slide, f"TimelineText{i+1}", Inches(1 + i*3.6 - 1), Inches(4), Inches(2.6), Inches(1.5), 
                     f"TimelineText{i+1} Placeholder", 12)

# Layout 11: Title with Bulleted List
def add_layout_11_title_bulleted_list():
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, "TitleBulletedListTitle", Inches(1), Inches(0.5), Inches(11.33), Inches(1), 
                 "TitleBulletedListTitle Placeholder", 32, True, PP_ALIGN.CENTER)
    add_text_box(slide, "TitleBulletedListContent", Inches(1), Inches(2), Inches(11.33), Inches(5), 
                 "TitleBulletedListContent Placeholder", 18, False, PP_ALIGN.LEFT)

# Layout 12: Image Gallery (3x2 Grid)
def add_layout_12_image_gallery():
    slide = prs.slides.add_slide(blank_layout)
    for row in range(3):
        for col in range(2):
            add_image_placeholder(slide, f"ImageGalleryImage{row*2+col+1}", 
                                 Inches(0.5 + col*6.5), Inches(0.5 + row*2.3), Inches(6), Inches(1.5))
            add_text_box(slide, f"ImageGalleryCaption{row*2+col+1}", 
                         Inches(0.5 + col*6.5), Inches(2.1 + row*2.3), Inches(6), Inches(0.5), 
                         f"ImageGalleryCaption{row*2+col+1} Placeholder", 12, False, PP_ALIGN.CENTER)

# Layout 13: Process Flow (4 Steps)
def add_layout_13_process_flow():
    slide = prs.slides.add_slide(blank_layout)
    for i in range(4):
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1 + i*3.3), Inches(1), Inches(0.5), Inches(0.5))
        marker.name = f"ProcessFlowMarker{i+1}"
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                                          Inches(1.5 + i*3.3), Inches(1.15), Inches(2.8), Inches(0.2))
            arrow.name = f"ProcessFlowArrow{i+1}"
        add_text_box(slide, f"ProcessFlowTitle{i+1}", Inches(1 + i*3.3 - 0.5), Inches(1.7), Inches(3.3), Inches(0.5), 
                     f"ProcessFlowTitle{i+1} Placeholder", 16, True, PP_ALIGN.CENTER)
        add_text_box(slide, f"ProcessFlowDescription{i+1}", Inches(1 + i*3.3 - 0.5), Inches(2.3), Inches(3.3), Inches(2), 
                     f"ProcessFlowDescription{i+1} Placeholder", 12, False, PP_ALIGN.CENTER)

# Layout 14: Team Introduction
def add_layout_14_team_introduction():
    slide = prs.slides.add_slide(blank_layout)
    add_text_box(slide, "TeamIntroductionTitle", Inches(1), Inches(0.5), Inches(11.33), Inches(1), 
                 "TeamIntroductionTitle Placeholder", 32, True, PP_ALIGN.CENTER)
    for i in range(3):
        add_image_placeholder(slide, f"TeamIntroductionImage{i+1}", 
                             Inches(1 + i*4.3), Inches(2), Inches(2), Inches(2))
        add_text_box(slide, f"TeamIntroductionName{i+1}", 
                     Inches(1 + i*4.3), Inches(4.2), Inches(2), Inches(0.5), 
                     f"TeamIntroductionName{i+1} Placeholder", 16, True, PP_ALIGN.CENTER)
        add_text_box(slide, f"TeamIntroductionRole{i+1}", 
                     Inches(1 + i*4.3), Inches(4.8), Inches(2), Inches(0.5), 
                     f"TeamIntroductionRole{i+1} Placeholder", 12, False, PP_ALIGN.CENTER)

# Generate the template file on startup
def create_template():
    add_layout_1_title_slide()
    add_layout_2_2x2_grid()
    add_layout_3_image_caption_content()
    add_layout_4_image_two_sections()
    add_layout_5_full_image_overlay()
    add_layout_6_three_columns()
    add_layout_7_title_two_images()
    add_layout_8_comparison()
    add_layout_9_quote()
    add_layout_10_timeline()
    add_layout_11_title_bulleted_list()
    add_layout_12_image_gallery()
    add_layout_13_process_flow()
    add_layout_14_team_introduction()
    prs.save("all_layouts_template.potx")

# Map layout names to indices (0-based)
layout_map = {
    "title_slide": 0,
    "2x2_grid": 1,
    "image_caption_content": 2,
    "image_two_sections": 3,
    "full_image_overlay": 4,
    "three_columns": 5,
    "title_two_images": 6,
    "comparison": 7,
    "quote": 8,
    "timeline": 9,
    "title_bulleted_list": 10,
    "image_gallery": 11,
    "process_flow": 12,
    "team_introduction": 13
}

# Helper function to convert hex color to RGB
def hex_to_rgb(hex_color):
    hex_color = hex_color.replace("0x", "")
    return RGBColor(int(hex_color[2:4], 16), int(hex_color[4:6], 16), int(hex_color[6:8], 16))

# Helper function to map alignment
def get_alignment(alignment_x):
    if alignment_x == -1:
        return PP_ALIGN.LEFT
    elif alignment_x == 1:
        return PP_ALIGN.RIGHT
    return PP_ALIGN.CENTER

# Helper function to download and validate image
def download_image(url, filename):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        content = response.content
        image = Image.open(io.BytesIO(content))
        image.verify()
        image.close()
        with open(filename, "wb") as f:
            f.write(content)
        return filename
    except (requests.RequestException, Image.UnidentifiedImageError, Exception) as e:
        print(f"Failed to download or validate image from {url}: {str(e)}")
        return None

# Function to generate PPT based on layout and JSON data
def generate_ppt(layout_identifier, json_data):
    # Load the template
    prs = Presentation("all_layouts_template.potx")

    # Determine the layout index
    if isinstance(layout_identifier, int):
        layout_idx = layout_identifier - 1  # Convert to 0-based index
    elif isinstance(layout_identifier, str):
        layout_idx = layout_map.get(layout_identifier.lower())
        if layout_idx is None:
            return None, f"Invalid layout name: {layout_identifier}. Available layouts: {list(layout_map.keys())}"
    else:
        return None, "Layout identifier must be an integer (1-14) or a valid layout name."

    if not (0 <= layout_idx < 14):
        return None, "Layout number must be between 1 and 14."

    # Create a new slide using the selected layout
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set a black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # Use the corresponding template slide
    template_slide = prs.slides[layout_idx]

    # Copy shapes from the template slide
    for template_shape in template_slide.shapes:
        if template_shape.has_text_frame and template_shape.name != "TimelineLine":
            new_shape = slide.shapes.add_textbox(
                left=template_shape.left,
                top=template_shape.top,
                width=template_shape.width,
                height=template_shape.height
            )
            new_shape.name = template_shape.name
            new_shape.text = template_shape.text
            new_shape.text_frame.word_wrap = True
        else:
            new_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE if "Marker" not in template_shape.name and "Arrow" not in template_shape.name else MSO_SHAPE.OVAL if "Marker" in template_shape.name else MSO_SHAPE.RIGHT_ARROW,
                left=template_shape.left,
                top=template_shape.top,
                width=template_shape.width,
                height=template_shape.height
            )
            new_shape.name = template_shape.name

    # Process sections from JSON data
    for idx, section in enumerate(json_data["sections"]):
        placeholder_name = section.get("placeholderName")
        shape = None
        for s in slide.shapes:
            if s.name == placeholder_name:
                shape = s
                break

        if not shape:
            print(f"Shape {placeholder_name} not found in layout {layout_idx + 1}!")
            continue

        if section["type"] == "content_text":
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.text = section["data"].replace("\n", "\r\n")
            text_frame.word_wrap = True
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = get_alignment(section["alignmentX"])
                for run in paragraph.runs:
                    run.font.size = Pt(section["fontSize"])
                    run.font.bold = section["isBold"]
                    run.font.color.rgb = hex_to_rgb(section["colorHex"])
        elif section["type"] == "image":
            image_url = section.get("imageUrl")
            if image_url:
                image_filename = f"temp_image_{idx}.png"
                image_path = download_image(image_url, image_filename)
                if image_path:
                    try:
                        slide.shapes.add_picture(
                            image_path,
                            left=shape.left,
                            top=shape.top,
                            width=shape.width,
                            height=shape.height
                        )
                        shape._element.getparent().remove(shape._element)
                    finally:
                        os.remove(image_path)
                else:
                    shape.text = f"Image failed to load: {image_url}"

    # Add speaker notes
    if json_data.get("speakerNotes"):
        slide.notes_slide.notes_text_frame.text = json_data["speakerNotes"]

    # Save the presentation to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
        prs.save(temp_file.name)
        temp_file_path = temp_file.name

    return temp_file_path, None

# API endpoint to generate and return PPTX file
@app.route('/generate_ppt', methods=['POST'])
def generate_ppt_endpoint():
    try:
        # Check if request contains JSON data
        if not request.is_json:
            return jsonify({"error": "Request must contain JSON data"}), 400

        data = request.get_json()

        # Validate required fields
        if "template_name" not in data:
            return jsonify({"error": "Missing 'template_name' in request body"}), 400
        if "slides_data" not in data:
            return jsonify({"error": "Missing 'slides_data' in request body"}), 400

        template_name = data["template_name"]
        slides_data = data["slides_data"]

        # Generate the PPT
        pptx_path, error = generate_ppt(template_name, slides_data)

        if error:
            return jsonify({"error": error}), 400

        # Send the file as a response
        response = send_file(
            pptx_path,
            as_attachment=True,
            download_name=f"presentation_{template_name}.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        # Clean up the temporary file after sending
        os.remove(pptx_path)

        return response

    except Exception as e:
        return jsonify({"error": f"Failed to generate PPTX: {str(e)}"}), 500

# Create the template on startup
if __name__ == "__main__":
    create_template()
    app.run(host="0.0.0.0", port=5000, debug=True)