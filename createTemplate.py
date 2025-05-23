from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # Blank layout

# Helper function to add a text box with formatting
def add_text_box(slide, name, left, top, width, height, text, font_size, is_bold=False, alignment=PP_ALIGN.LEFT):
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_box.name = name
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = is_bold

# Helper function to add an image placeholder
def add_image_placeholder(slide, name, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.name = name
    shape.text = f"{name} Placeholder"
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(14)

# Layout 1: Title Slide
def add_layout_1_title_slide():
    slide = prs.slides.add_slide(blank_layout)
    # Title
    add_text_box(slide, "TitleSlideTitle", Inches(1), Inches(1), Inches(11.33), Inches(2), 
                 "TitleSlide Title Placeholder", 44, True, PP_ALIGN.CENTER)
    # Subtitle
    add_text_box(slide, "TitleSlideSubtitle", Inches(1), Inches(3.5), Inches(11.33), Inches(1), 
                 "TitleSlide Subtitle Placeholder", 24, False, PP_ALIGN.CENTER)
    # Logo
    add_image_placeholder(slide, "TitleSlideLogo", Inches(5.665), Inches(5.5), Inches(2), Inches(1))

# Layout 2: 2x2 Grid Layout
def add_layout_2_2x2_grid():
    slide = prs.slides.add_slide(blank_layout)
    # Top-left: Title + Content
    add_text_box(slide, "GridTopLeftTitle", Inches(0.5), Inches(0.3), Inches(6), Inches(0.5), 
                 "GridTopLeftTitle Placeholder", 18, True)
    add_text_box(slide, "GridTopLeftContent", Inches(0.5), Inches(0.9), Inches(6), Inches(2.6), 
                 "GridTopLeftContent Placeholder", 14)
    # Top-right: Image
    add_image_placeholder(slide, "GridTopRightImage", Inches(6.83), Inches(0.5), Inches(6), Inches(3))
    # Bottom-left: Image
    add_image_placeholder(slide, "GridBottomLeftImage", Inches(0.5), Inches(4), Inches(6), Inches(3))
    # Bottom-right: Title + Content
    add_text_box(slide, "GridBottomRightTitle", Inches(6.83), Inches(3.8), Inches(6), Inches(0.5), 
                 "GridBottomRightTitle Placeholder", 18, True)
    add_text_box(slide, "GridBottomRightContent", Inches(6.83), Inches(4.4), Inches(6), Inches(2.6), 
                 "GridBottomRightContent Placeholder", 14)

# Layout 3: Image with Caption and Content
def add_layout_3_image_caption_content():
    slide = prs.slides.add_slide(blank_layout)
    # Left: Image + Caption
    add_image_placeholder(slide, "ImageCaptionLeftImage", Inches(0.5), Inches(0.5), Inches(6), Inches(5.5))
    add_text_box(slide, "ImageCaptionLeftCaption", Inches(0.5), Inches(6.1), Inches(6), Inches(1), 
                 "ImageCaptionLeftCaption Placeholder", 12)
    # Right: Title + Content
    add_text_box(slide, "ImageCaptionRightTitle", Inches(6.83), Inches(0.5), Inches(6), Inches(0.5), 
                 "ImageCaptionRightTitle Placeholder", 18, True)
    add_text_box(slide, "ImageCaptionRightContent", Inches(6.83), Inches(1.1), Inches(6), Inches(6.1), 
                 "ImageCaptionRightContent Placeholder", 14)

# Layout 4: Image Left, Two Sections Right
def add_layout_4_image_two_sections():
    slide = prs.slides.add_slide(blank_layout)
    # Left: Image
    add_image_placeholder(slide, "ImageTwoSectionsLeftImage", Inches(0.5), Inches(0.5), Inches(6), Inches(6.5))
    # Right: Two Sections
    add_text_box(slide, "ImageTwoSectionsRightTitle1", Inches(6.83), Inches(0.5), Inches(6), Inches(0.5), 
                 "ImageTwoSectionsRightTitle1 Placeholder", 18, True)
    add_text_box(slide, "ImageTwoSectionsRightContent1", Inches(6.83), Inches(1.1), Inches(6), Inches(2.9), 
                 "ImageTwoSectionsRightContent1 Placeholder", 14)
    add_text_box(slide, "ImageTwoSectionsRightTitle2", Inches(6.83), Inches(4.2), Inches(6), Inches(0.5), 
                 "ImageTwoSectionsRightTitle2 Placeholder", 18, True)
    add_text_box(slide, "ImageTwoSectionsRightContent2", Inches(6.83), Inches(4.8), Inches(6), Inches(2.2), 
                 "ImageTwoSectionsRightContent2 Placeholder", 14)

# Layout 5: Full Slide Image with Overlay Text
def add_layout_5_full_image_overlay():
    slide = prs.slides.add_slide(blank_layout)
    # Full-slide Image
    add_image_placeholder(slide, "FullImageOverlayImage", Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    # Overlay Text
    add_text_box(slide, "FullImageOverlayText", Inches(1), Inches(3), Inches(11.33), Inches(1.5), 
                 "FullImageOverlayText Placeholder", 32, True, PP_ALIGN.CENTER)

# Layout 6: Three-Column Layout
def add_layout_6_three_columns():
    slide = prs.slides.add_slide(blank_layout)
    # Column 1
    add_text_box(slide, "ThreeColumnsTitle1", Inches(0.5), Inches(0.5), Inches(4), Inches(0.5), 
                 "ThreeColumnsTitle1 Placeholder", 18, True)
    add_text_box(slide, "ThreeColumnsContent1", Inches(0.5), Inches(1.1), Inches(4), Inches(6.1), 
                 "ThreeColumnsContent1 Placeholder", 14)
    # Column 2
    add_text_box(slide, "ThreeColumnsTitle2", Inches(4.66), Inches(0.5), Inches(4), Inches(0.5), 
                 "ThreeColumnsTitle2 Placeholder", 18, True)
    add_text_box(slide, "ThreeColumnsContent2", Inches(4.66), Inches(1.1), Inches(4), Inches(6.1), 
                 "ThreeColumnsContent2 Placeholder", 14)
    # Column 3
    add_text_box(slide, "ThreeColumnsTitle3", Inches(8.83), Inches(0.5), Inches(4), Inches(0.5), 
                 "ThreeColumnsTitle3 Placeholder", 18, True)
    add_text_box(slide, "ThreeColumnsContent3", Inches(8.83), Inches(1.1), Inches(4), Inches(6.1), 
                 "ThreeColumnsContent3 Placeholder", 14)

# Layout 7: Title with Two Images Below
def add_layout_7_title_two_images():
    slide = prs.slides.add_slide(blank_layout)
    # Title
    add_text_box(slide, "TitleTwoImagesTitle", Inches(1), Inches(0.5), Inches(11.33), Inches(1), 
                 "TitleTwoImagesTitle Placeholder", 32, True, PP_ALIGN.CENTER)
    # Images
    add_image_placeholder(slide, "TitleTwoImagesLeftImage", Inches(0.5), Inches(2), Inches(6), Inches(4))
    add_image_placeholder(slide, "TitleTwoImagesRightImage", Inches(6.83), Inches(2), Inches(6), Inches(4))

# Layout 8: Comparison Layout
def add_layout_8_comparison():
    slide = prs.slides.add_slide(blank_layout)
    # Left: Title + Content
    add_text_box(slide, "ComparisonLeftTitle", Inches(0.5), Inches(0.5), Inches(6), Inches(0.5), 
                 "ComparisonLeftTitle Placeholder", 18, True)
    add_text_box(slide, "ComparisonLeftContent", Inches(0.5), Inches(1.1), Inches(6), Inches(6.1), 
                 "ComparisonLeftContent Placeholder", 14)
    # Right: Title + Content
    add_text_box(slide, "ComparisonRightTitle", Inches(6.83), Inches(0.5), Inches(6), Inches(0.5), 
                 "ComparisonRightTitle Placeholder", 18, True)
    add_text_box(slide, "ComparisonRightContent", Inches(6.83), Inches(1.1), Inches(6), Inches(6.1), 
                 "ComparisonRightContent Placeholder", 14)

# Layout 9: Quote Slide
def add_layout_9_quote():
    slide = prs.slides.add_slide(blank_layout)
    # Quote
    add_text_box(slide, "QuoteText", Inches(1), Inches(1.5), Inches(11.33), Inches(3), 
                 "QuoteText Placeholder", 28, False, PP_ALIGN.CENTER)
    # Attribution
    add_text_box(slide, "QuoteAttribution", Inches(1), Inches(5), Inches(11.33), Inches(1), 
                 "QuoteAttribution Placeholder", 18, False, PP_ALIGN.CENTER)

# Layout 10: Timeline Layout
def add_layout_10_timeline():
    slide = prs.slides.add_slide(blank_layout)
    # Timeline Line (represented as a thin rectangle)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.5), Inches(11.33), Inches(0.1))
    shape.name = "TimelineLine"
    # Milestones
    for i in range(4):
        # Milestone Marker (small circle)
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1 + i*3.6), Inches(3.4), Inches(0.3), Inches(0.3))
        marker.name = f"TimelineMarker{i+1}"
        # Milestone Text
        add_text_box(slide, f"TimelineText{i+1}", Inches(1 + i*3.6 - 1), Inches(4), Inches(2.6), Inches(1.5), 
                     f"TimelineText{i+1} Placeholder", 12)

# Layout 11: Title with Bulleted List
def add_layout_11_title_bulleted_list():
    slide = prs.slides.add_slide(blank_layout)
    # Title
    add_text_box(slide, "TitleBulletedListTitle", Inches(1), Inches(0.5), Inches(11.33), Inches(1), 
                 "TitleBulletedListTitle Placeholder", 32, True, PP_ALIGN.CENTER)
    # Bulleted List
    add_text_box(slide, "TitleBulletedListContent", Inches(1), Inches(2), Inches(11.33), Inches(5), 
                 "TitleBulletedListContent Placeholder", 18, False, PP_ALIGN.LEFT)

# Layout 12: Image Gallery (3x2 Grid)
def add_layout_12_image_gallery():
    slide = prs.slides.add_slide(blank_layout)
    for row in range(3):
        for col in range(2):
            # Image Placeholder
            add_image_placeholder(slide, f"ImageGalleryImage{row*2+col+1}", 
                                 Inches(0.5 + col*6.5), Inches(0.5 + row*2.3), Inches(6), Inches(1.5))
            # Caption
            add_text_box(slide, f"ImageGalleryCaption{row*2+col+1}", 
                         Inches(0.5 + col*6.5), Inches(2.1 + row*2.3), Inches(6), Inches(0.5), 
                         f"ImageGalleryCaption{row*2+col+1} Placeholder", 12, False, PP_ALIGN.CENTER)

# Layout 13: Process Flow (4 Steps)
def add_layout_13_process_flow():
    slide = prs.slides.add_slide(blank_layout)
    for i in range(4):
        # Step Marker (circle)
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1 + i*3.3), Inches(1), Inches(0.5), Inches(0.5))
        marker.name = f"ProcessFlowMarker{i+1}"
        # Arrow between steps (except after the last step)
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                                          Inches(1.5 + i*3.3), Inches(1.15), Inches(2.8), Inches(0.2))
            arrow.name = f"ProcessFlowArrow{i+1}"
        # Step Title
        add_text_box(slide, f"ProcessFlowTitle{i+1}", Inches(1 + i*3.3 - 0.5), Inches(1.7), Inches(3.3), Inches(0.5), 
                     f"ProcessFlowTitle{i+1} Placeholder", 16, True, PP_ALIGN.CENTER)
        # Step Description
        add_text_box(slide, f"ProcessFlowDescription{i+1}", Inches(1 + i*3.3 - 0.5), Inches(2.3), Inches(3.3), Inches(2), 
                     f"ProcessFlowDescription{i+1} Placeholder", 12, False, PP_ALIGN.CENTER)

# Layout 14: Team Introduction
def add_layout_14_team_introduction():
    slide = prs.slides.add_slide(blank_layout)
    # Title
    add_text_box(slide, "TeamIntroductionTitle", Inches(1), Inches(0.5), Inches(11.33), Inches(1), 
                 "TeamIntroductionTitle Placeholder", 32, True, PP_ALIGN.CENTER)
    # Team Members
    for i in range(3):
        # Image Placeholder
        add_image_placeholder(slide, f"TeamIntroductionImage{i+1}", 
                             Inches(1 + i*4.3), Inches(2), Inches(2), Inches(2))
        # Name
        add_text_box(slide, f"TeamIntroductionName{i+1}", 
                     Inches(1 + i*4.3), Inches(4.2), Inches(2), Inches(0.5), 
                     f"TeamIntroductionName{i+1} Placeholder", 16, True, PP_ALIGN.CENTER)
        # Role
        add_text_box(slide, f"TeamIntroductionRole{i+1}", 
                     Inches(1 + i*4.3), Inches(4.8), Inches(2), Inches(0.5), 
                     f"TeamIntroductionRole{i+1} Placeholder", 12, False, PP_ALIGN.CENTER)

# Add all layouts
add_layout_1_title_slide()
add_layout_2_2x2_grid()
add_layout_3_image_caption_content()
add_layout_4_image_two_sections()
add_layout_5_full_image_overlay()
add_layout_6_three_columns()
add_layout_7_title_two_images()
add_layout_8_comparison()
add_layout_9_quote()
add_layout_10_timeline()
add_layout_11_title_bulleted_list()
add_layout_12_image_gallery()
add_layout_13_process_flow()
add_layout_14_team_introduction()

# Save the template
prs.save("all_layouts_template.potx")
print("Template with 14 layouts created successfully!")